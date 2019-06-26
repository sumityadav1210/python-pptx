files_array = ['/home/sumit/Desktop/mergeppt/static/pptsources/RU/industries/auto/trends/L2/trends-L2.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/industries/auto/problems/L1/problems-L1.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/industries/auto/howsolved/L0/solved-L0.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/customerrefs/auto/global/is/templ.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/customerrefs/auto/regional/ci/templ.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/customerrefs/auto/local/is/templ.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/personalrefs/templ2.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/personalrefs/templ.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/solutiondetails/car/L0/templ.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/solutiondetails/platform/L2/templ.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/addendums/customer_success/success.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/addendums/platform_integration/integration.pptx', '/home/sumit/Desktop/mergeppt/static/pptsources/RU/thankyou/typ.pptx']

from django.shortcuts import render
from django.http import HttpResponse,HttpResponseRedirect
import os,ast
import pptx
from pptx import Presentation
from pptx.parts.chart import ChartPart
from pptx import Presentation
import six
import copy

# Create your views here.


base_url = "/home/sumit/Desktop/mergeppt/static/pptsources"

def _get_blank_slide_layout(pres):
	layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
	min_items = min(layout_items_count)
	blank_layout_id = layout_items_count.index(min_items)
	return pres.slide_layouts[blank_layout_id]


def copy_slide(pres1,pres2,index):
	source = pres1.slides[index]
	blank_slide_layout = _get_blank_slide_layout(pres1)
	dest = pres2.slides.add_slide(blank_slide_layout)
	for shp in source.shapes:
		el = shp.element
		newel = copy.deepcopy(el)
		dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
	
	for key, value in six.iteritems(source.part.rels):
		# Make sure we don't copy a notesSlide relation as that won't exist
		if not "notesSlide" in value.reltype:
			target = value._target
			print(value.reltype, value._target, value.rId)
			# if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
			if "chart" in value.reltype:
				partname = target.package.next_partname(ChartPart.partname_template)
				xlsx_blob = target.chart_workbook.xlsx_part.blob
				target = ChartPart(partname, target.content_type,copy.deepcopy(target._element), package=target.package)
				target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(xlsx_blob, target.package)
			
			dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
	
	return dest


prs2 = Presentation(files_array.pop(0))

for file1 in files_array:
	prs1 = Presentation(file1)
	i=0
	for slide in prs1.slides:
		copy_slide(prs1, prs2, i)
		i=i+1
	
prs2.save("test.pptx")





